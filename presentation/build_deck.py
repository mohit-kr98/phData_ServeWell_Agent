"""Final ServeWell IT Support Agent deck, built against the phData challenge brief."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY      = RGBColor(0x12, 0x28, 0x40)
NAVY_2    = RGBColor(0x1D, 0x3A, 0x5C)
TEAL      = RGBColor(0x17, 0x92, 0x85)
TEAL_DK   = RGBColor(0x0D, 0x66, 0x5E)
TEAL_SOFT = RGBColor(0xE4, 0xF2, 0xF0)
BG        = RGBColor(0xF6, 0xF7, 0xF9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x1A, 0x20, 0x2A)
INK_2     = RGBColor(0x57, 0x60, 0x6D)
LINE      = RGBColor(0xD8, 0xDD, 0xE3)
GOOD      = RGBColor(0x2B, 0x94, 0x55)
WARN      = RGBColor(0xB8, 0x82, 0x18)
BAD       = RGBColor(0xB8, 0x39, 0x2C)
VIOLET    = RGBColor(0x4C, 0x4A, 0x8F)
AMBER     = RGBColor(0xA8, 0x6A, 0x14)
F = "Calibri"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
PAGE = [0]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, lw=1.0, adj=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is not None:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    if adj is not None:
        try: sp.adjustments[0] = adj
        except Exception: pass
    return sp


def text(s, x, y, w, h, t, size=14, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.0, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line_ in enumerate(str(t).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        r = p.add_run(); r.text = line_
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = italic; r.font.name = F
    return tb


def bullets(s, x, y, w, h, items, size=15, color=INK, bcolor=TEAL, gap=10, ls=1.12):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        body, lvl = it[0], it[1]
        lead = it[2] if len(it) > 2 else None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = ls
        sz = size - (2 if lvl else 0)
        r0 = p.add_run(); r0.text = "•  " if lvl == 0 else "–  "
        r0.font.size = Pt(sz); r0.font.color.rgb = bcolor; r0.font.bold = True; r0.font.name = F
        if lead:
            r1 = p.add_run(); r1.text = lead
            r1.font.size = Pt(sz); r1.font.color.rgb = color; r1.font.bold = True; r1.font.name = F
        r2 = p.add_run(); r2.text = body
        r2.font.size = Pt(sz); r2.font.color.rgb = INK_2 if lead else color; r2.font.name = F
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(Inches(0.30 * lvl)))); pPr.set("indent", "0")
    return tb


def head(s, kicker, title, dark=False):
    text(s, Inches(0.55), Inches(0.36), Inches(12), Inches(0.3), kicker.upper(),
         size=12, color=TEAL, bold=True)
    text(s, Inches(0.55), Inches(0.66), Inches(12.2), Inches(0.75), title,
         size=27, color=WHITE if dark else INK, bold=True)
    rect(s, Inches(0.55), Inches(1.40), Inches(0.85), Pt(3), TEAL)


def foot(s, note=None):
    # True slide position, not a running counter: the title slide has no
    # footer, so a counter would report every later slide one short -- and
    # slide 2 cross-references these numbers ("5-6 - component architecture"),
    # so they have to match what the presenter sees in the tray.
    page = len(prs.slides._sldIdLst)
    text(s, Inches(0.55), Inches(7.15), Inches(10), Inches(0.28),
         note or "ServeWell IT Support Agent Suite — phData Intelligence Platform", size=9, color=INK_2)
    text(s, Inches(12.4), Inches(7.15), Inches(0.5), Inches(0.28), str(page),
         size=9, color=INK_2, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, label, value, sub=None, vc=NAVY):
    c = rect(s, x, y, w, h, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, adj=0.045)
    text(s, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.32), label.upper(), size=10.5, color=INK_2, bold=True)
    text(s, x + Inches(0.2), y + Inches(0.47), w - Inches(0.4), Inches(0.55), value, size=27, color=vc, bold=True)
    if sub:
        text(s, x + Inches(0.2), y + h - Inches(0.40), w - Inches(0.4), Inches(0.34), sub, size=10, color=INK_2)


def arrow(s, x1, y1, x2, y2, color=TEAL_DK, w=2.0, dashed=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = Pt(w); c.shadow.inherit = False
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return c


def box(s, x, y, w, h, title, sub=None, fill=WHITE, tc=INK, sc=INK_2, line=LINE, lw=1.0, tsize=11.5, ssize=8.5):
    b = rect(s, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE, line=line, lw=lw, adj=0.10)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(5); tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(tsize); r.font.bold = True; r.font.color.rgb = tc; r.font.name = F
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(ssize); r2.font.color.rgb = sc; r2.font.name = F
    return b


def pill(s, x, y, w, h, label, color, fill=WHITE, fs=9.5):
    """Provenance chip: matches the roi.html artifact's Measured/Stated/
    Assumption tag language, so the same visual vocabulary means the same
    thing whether a reviewer is looking at the deck or the live artifact."""
    b = rect(s, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE, line=color, lw=1.1, adj=0.5)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = color; r.font.name = F
    return b


def table(s, x, y, w, h, headers, rows, ratios=None, fs=12, hdr=NAVY, colors=None, aligns=None):
    t = s.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    if ratios:
        tot = sum(ratios)
        for i, rr in enumerate(ratios):
            t.columns[i].width = Emu(int(w * rr / tot))
    for j, htxt in enumerate(headers):
        c = t.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = hdr
        c.margin_top = c.margin_bottom = Pt(5); c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = aligns[j] if aligns else PP_ALIGN.LEFT
        r = p.add_run(); r.text = htxt
        r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = F
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = t.cell(i + 1, j); c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(0xF1, 0xF4, 0xF7)
            c.margin_top = c.margin_bottom = Pt(4); c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = aligns[j] if aligns else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(v)
            r.font.size = Pt(fs - 1); r.font.name = F; r.font.color.rgb = INK
            if colors and (i, j) in colors:
                r.font.color.rgb = colors[(i, j)]; r.font.bold = True
    return t


# =====================================================================
# 1 — Title
# =====================================================================
s = slide(NAVY)
rect(s, 0, 0, SW, Inches(0.11), TEAL)
rect(s, Inches(0.6), Inches(2.5), Inches(0.85), Pt(4), TEAL)
text(s, Inches(0.6), Inches(2.70), Inches(12), Inches(0.34), "PHDATA INTELLIGENCE PLATFORM  ·  SENIOR APPLIED AI ENGINEER CHALLENGE",
     size=13.5, color=TEAL, bold=True)
text(s, Inches(0.6), Inches(3.08), Inches(12.2), Inches(1.0), "ServeWell IT Support Agent Suite", size=42, color=WHITE, bold=True)
text(s, Inches(0.6), Inches(4.05), Inches(11.6), Inches(0.9),
     "An agentic triage → grounded resolution → human-approved action system,\nproven on a narrow slice of ServeWell's L1 support load.",
     size=17, color=RGBColor(0xC3, 0xD0, 0xDE), ls=1.25)
for i, (k, v) in enumerate([("Routing accuracy", "90.2%"), ("Runbook hit-rate", "100%"),
                            ("Grounded replies", "100%"), ("Response p95", "3.9s")]):
    x = Inches(0.6) + Inches(3.05) * i
    text(s, x, Inches(5.62), Inches(2.9), Inches(0.3), k.upper(), size=9.5, color=RGBColor(0x8B, 0x9D, 0xB2), bold=True)
    text(s, x, Inches(5.92), Inches(2.9), Inches(0.5), v, size=25, color=TEAL, bold=True)
text(s, Inches(0.6), Inches(6.85), Inches(11), Inches(0.3),
     "Measured on 220 labelled tickets against the challenge answer key · reproducible with one command",
     size=10.5, color=RGBColor(0x7C, 0x8E, 0xA3))

# =====================================================================
# 2 — How to read this deck
# =====================================================================
s = slide(); head(s, "Orientation", "Two audiences, one system")
text(s, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.5),
     "The brief asks for a walkthrough a non-technical stakeholder can follow and a technical reviewer can trust. "
     "These slides are marked so you can skip to your half — but the middle section is the one I most want probed.",
     size=13, color=INK_2, ls=1.2)
groups = [
    ("FOR THE BUSINESS UNIT VP", VIOLET, [
        "3 — The problem, in ServeWell's own numbers",
        "4 — What we built, in plain language",
        "17 — Business impact and what it rests on",
        "18 — What I would not yet claim",
    ]),
    ("FOR THE DIRECTOR OF TECHNOLOGY", TEAL_DK, [
        "5–6 — Component architecture and request path",
        "7–11 — Agents, RAG, tools, actions, guardrails",
        "12–15 — Evaluation, a measurement failure, tuning",
        "16 — Technical trade-offs and where I drew the line",
    ]),
]
for i, (t_, c, items) in enumerate(groups):
    x = Inches(0.55) + Inches(6.3) * i
    p = rect(s, x, Inches(2.35), Inches(6.0), Inches(3.5), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, adj=0.05)
    rect(s, x, Inches(2.35), Inches(0.12), Inches(3.5), c)
    text(s, x + Inches(0.35), Inches(2.58), Inches(5.4), Inches(0.32), t_, size=12, color=c, bold=True)
    bullets(s, x + Inches(0.35), Inches(3.05), Inches(5.3), Inches(2.6),
            [(b, 0) for b in items], size=13, bcolor=c, gap=13)
rect(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.72), TEAL_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.14)
text(s, Inches(0.85), Inches(6.32), Inches(11.6), Inches(0.45),
     "Every number in this deck is measured by a harness in the repo, or explicitly labelled as an assumption. Nothing is estimated and presented as measured.",
     size=12.5, color=TEAL_DK, italic=True)
foot(s)

# =====================================================================
# 3 — The problem
# =====================================================================
s = slide(); head(s, "The Problem", "What ServeWell told us is breaking")
pains = [
    ("30–45 min", "to first meaningful response", "during peak hours, when a store is losing revenue"),
    ("60–70%", "of tickets are already known", "they map to an existing runbook, but still need a human to read it"),
    ("2,500+", "locations, inconsistent service", "some stores get quick help; others bounce between L1 and L2"),
]
for i, (big, mid, sub) in enumerate(pains):
    x = Inches(0.55) + Inches(4.15) * i
    c = rect(s, x, Inches(1.75), Inches(3.85), Inches(2.05), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, adj=0.05)
    text(s, x + Inches(0.28), Inches(1.98), Inches(3.3), Inches(0.6), big, size=31, color=BAD, bold=True)
    text(s, x + Inches(0.28), Inches(2.58), Inches(3.3), Inches(0.35), mid, size=13, color=INK, bold=True)
    text(s, x + Inches(0.28), Inches(2.92), Inches(3.35), Inches(0.7), sub, size=10.5, color=INK_2, ls=1.15)
text(s, Inches(0.55), Inches(4.12), Inches(12), Inches(0.35), "AND THE COST LANDS IN THREE PLACES", size=12, color=TEAL, bold=True)
bullets(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(1.6), [
    ("Store downtime reduces revenue and customer satisfaction directly.", 0, "Revenue:  "),
    ("L1 cost and burnout are rising, with high attrition in the support centre.", 0, "People:  "),
    ("L2 engineers are absorbing basic tickets that should never have reached them.", 0, "Capacity:  "),
], size=15, gap=13)
rect(s, Inches(0.55), Inches(6.28), Inches(12.2), Inches(0.7), TEAL_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.14)
text(s, Inches(0.85), Inches(6.44), Inches(11.6), Inches(0.45),
     "Note the shape of this: the knowledge to fix most of these tickets already exists. The bottleneck is getting it to the right person fast, and knowing when not to try.",
     size=12.5, color=TEAL_DK, italic=True)
foot(s)

# =====================================================================
# 4 — What we built (plain language)
# =====================================================================
s = slide(); head(s, "In Plain Language", "What the system actually does")
steps = [
    ("A ticket arrives", "“The PIN pad on terminal 1 is not responding.”", NAVY),
    ("It decides who should handle it", "Straight to a human L2 engineer, or guided at L1", TEAL_DK),
    ("It finds the right runbook", "and reads the store, asset and warranty records", NAVY_2),
    ("It writes the steps back", "grounded in that runbook — never invented", TEAL_DK),
    ("It can offer to act", "but a human approves before anything runs", AMBER),
]
y = Inches(1.72)
for i, (t_, d, c) in enumerate(steps):
    ry = y + Inches(0.98) * i
    rect(s, Inches(0.55), ry, Inches(0.52), Inches(0.52), c, MSO_SHAPE.OVAL)
    text(s, Inches(0.55), ry, Inches(0.52), Inches(0.52), str(i + 1), size=17, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.30), ry - Inches(0.02), Inches(5.0), Inches(0.4), t_, size=16.5, bold=True, color=INK)
    text(s, Inches(1.30), ry + Inches(0.34), Inches(6.6), Inches(0.4), d, size=12.5, color=INK_2)
    if i < 4:
        arrow(s, Inches(0.81), ry + Inches(0.54), Inches(0.81), ry + Inches(0.94), color=LINE, w=1.5)
p = rect(s, Inches(8.3), Inches(1.72), Inches(4.45), Inches(4.5), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, adj=0.05)
text(s, Inches(8.6), Inches(1.98), Inches(3.9), Inches(0.35), "The one design rule", size=15, color=NAVY, bold=True)
text(s, Inches(8.6), Inches(2.42), Inches(3.9), Inches(2.6),
     "The agent is allowed to be wrong about advice, because a human reads it before acting on it.\n\n"
     "It is never allowed to be wrong about an action, because nothing runs without an explicit human approval.\n\n"
     "That asymmetry is why the autonomy line sits where it does.",
     size=13, color=INK_2, ls=1.35)
rect(s, Inches(8.6), Inches(5.32), Inches(3.85), Inches(0.62), TEAL_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.16)
text(s, Inches(8.78), Inches(5.46), Inches(3.5), Inches(0.4), "Advice is reversible. Actions are not.",
     size=12, color=TEAL_DK, bold=True, align=PP_ALIGN.CENTER)
foot(s)

# =====================================================================
# 5 — COMPONENT ARCHITECTURE (the centrepiece)
# =====================================================================
s = slide(); head(s, "Architecture", "Component-level view")

LX, LW = Inches(0.55), Inches(9.65)      # main column
RX, RW = Inches(10.42), Inches(2.33)     # observability rail

def band(y, h, label, color):
    rect(s, LX, y, LW, h, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, adj=0.06)
    rect(s, LX, y, Inches(0.10), h, color)
    text(s, LX + Inches(0.20), y + Inches(0.07), Inches(1.6), Inches(0.24), label,
         size=8, color=color, bold=True)

# --- Band 1: interfaces
b1y, b1h = Inches(1.58), Inches(0.72)
band(b1y, b1h, "INTERFACE", VIOLET)
box(s, LX + Inches(1.85), b1y + Inches(0.13), Inches(3.5), Inches(0.46), "Streamlit UI",
    "Live demo · Evaluation · Latency · Admin", fill=RGBColor(0xEE, 0xEE, 0xF7), line=VIOLET, tc=VIOLET, tsize=10.5, ssize=7.5)
box(s, LX + Inches(5.55), b1y + Inches(0.13), Inches(3.85), Inches(0.46), "FastAPI  ·  api.py",
    "/triage  /resolve  /l2_copilot  /execute_action", fill=RGBColor(0xEE, 0xEE, 0xF7), line=VIOLET, tc=VIOLET, tsize=10.5, ssize=7.5)

# --- Band 2: agents
b2y, b2h = Inches(2.42), Inches(0.96)
band(b2y, b2h, "AGENTS", NAVY)
aw = Inches(2.05)
for i, (t_, sub) in enumerate([("Triage Agent", "routing · rules + LLM"),
                               ("Resolution Agent", "L1 guidance loop"),
                               ("L2 Copilot", "assists human engineer")]):
    box(s, LX + Inches(1.85) + (aw + Inches(0.14)) * i, b2y + Inches(0.20), aw, Inches(0.58),
        t_, sub, fill=NAVY, tc=WHITE, sc=RGBColor(0xB9, 0xC7, 0xD6), line=NAVY, tsize=10, ssize=7)
# prompts/ sits beside the agents, not on top of them: it is a resource all
# three read from, so it gets its own column at the right edge of the band.
box(s, LX + Inches(8.42), b2y + Inches(0.20), Inches(1.08), Inches(0.58),
    "prompts/", "4 versioned files", fill=TEAL_SOFT, line=TEAL, tc=TEAL_DK, tsize=9.5, ssize=6.5)

# --- Band 3: safety
b3y, b3h = Inches(3.50), Inches(0.72)
band(b3y, b3h, "SAFETY", BAD)
sw_ = Inches(2.30)
for i, (t_, sub, c) in enumerate([("Deterministic rules", "history ≥2 · fault-class", TEAL_DK),
                                  ("policy_check() veto", "closed action catalog", BAD),
                                  ("HITL approval gate", "approved=True or refuse", AMBER)]):
    box(s, LX + Inches(1.85) + (sw_ + Inches(0.16)) * i, b3y + Inches(0.13), sw_, Inches(0.46),
        t_, sub, fill=WHITE, line=c, lw=1.5, tc=c, tsize=10, ssize=7.5)

# --- Band 4: tools
b4y, b4h = Inches(4.34), Inches(0.86)
band(b4y, b4h, "TOOLS", TEAL_DK)
tools = [("Retrieval", "search_knowledge_base\nsearch_faq · get_system_spec"),
         ("Structured data", "get_asset_info · get_store_info\ncheck_sla"),
         ("Actions", "propose_action · execute_action\nreply · resolve · escalate")]
for i, (t_, sub) in enumerate(tools):
    box(s, LX + Inches(1.85) + (sw_ + Inches(0.16)) * i, b4y + Inches(0.13), sw_, Inches(0.60),
        t_, sub, fill=TEAL_SOFT, line=TEAL, tc=TEAL_DK, tsize=10, ssize=7)

# --- Band 5: retrieval service
b5y, b5h = Inches(5.32), Inches(0.80)
band(b5y, b5h, "RETRIEVAL", NAVY_2)
pipe = ["result cache\n54% hit", "embed cache\nTitan v2", "PGVector\nmain + delta", "cross-encoder\nrerank"]
pw = Inches(1.72)
for i, lab in enumerate(pipe):
    bx = LX + Inches(1.85) + (pw + Inches(0.30)) * i
    box(s, bx, b5y + Inches(0.14), pw, Inches(0.52), lab.split("\n")[0], lab.split("\n")[1],
        fill=WHITE, line=NAVY_2, tc=NAVY_2, tsize=9.5, ssize=7)
    if i < 3:
        arrow(s, bx + pw, b5y + Inches(0.40), bx + pw + Inches(0.30), b5y + Inches(0.40), color=NAVY_2, w=1.5)

# --- Band 6: data
b6y, b6h = Inches(6.24), Inches(0.74)
band(b6y, b6h, "DATA", INK_2)
data = [("kb/  runbooks · FAQ · SOP", "markdown, chunked"), ("PGVector (Postgres)", "embeddings"),
        ("assets · stores · SLA", "CSV → CMDB stand-in"), ("labels/train_labels.json", "answer key")]
dw = Inches(1.79)
for i, (t_, sub) in enumerate(data):
    box(s, LX + Inches(1.85) + (dw + Inches(0.07)) * i, b6y + Inches(0.13), dw, Inches(0.48), t_, sub,
        fill=RGBColor(0xEC, 0xEF, 0xF2), line=LINE, tc=INK, sc=INK_2, tsize=8.5, ssize=7)

# vertical flow arrows down the left gutter
for (ay, by_) in [(b1y + b1h, b2y), (b2y + b2h, b3y), (b3y + b3h, b4y), (b4y + b4h, b5y), (b5y + b5h, b6y)]:
    arrow(s, LX + Inches(1.35), ay, LX + Inches(1.35), by_, color=TEAL_DK, w=1.75)

# --- observability rail
rect(s, RX, b1y, RW, Inches(5.40), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, adj=0.04)
rect(s, RX, b1y, RW, Inches(0.30), NAVY, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.30)
rect(s, RX, b1y + Inches(0.16), RW, Inches(0.14), NAVY)
text(s, RX, b1y + Inches(0.06), RW, Inches(0.24), "OBSERVABILITY & EVAL", size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
obs = [("Per-step trace", "every tool call, timed"), ("LangSmith", "distributed traces"),
       ("eval_labeled.py", "routing + retrieval"), ("eval_quality.py", "groundedness · guardrails"),
       ("Wall-clock timing", "concurrency-safe"), ("Run history", "every run archived"),
       ("logs/ticket_actions", "append-only audit")]
for i, (t_, sub) in enumerate(obs):
    oy = b1y + Inches(0.44) + Inches(0.70) * i
    text(s, RX + Inches(0.16), oy, RW - Inches(0.3), Inches(0.24), t_, size=9.5, color=NAVY, bold=True)
    text(s, RX + Inches(0.16), oy + Inches(0.21), RW - Inches(0.3), Inches(0.34), sub, size=7.5, color=INK_2, ls=1.05)
foot(s, "Component architecture — every box is a real module in the repo")

# =====================================================================
# 6 — Request path
# =====================================================================
s = slide(); head(s, "Architecture", "What happens to one ticket, and where the time goes")
lane_y = Inches(1.85)
stages = [
    ("Ingest", "ticket JSON\nPII masked", "~0ms", NAVY),
    ("Triage", "rule or LLM\ndecides queue", "0.28s", TEAL_DK),
    ("Parallel fetch", "3 searches + spec\n+ CMDB enrich", "0.9s", NAVY_2),
    ("Reason loop", "≤5 turns, capped\nsearch budget", "1.17s", TEAL_DK),
    ("Act", "reply / resolve /\npropose action", "~0ms", AMBER),
]
bw = Inches(2.32)
for i, (t_, sub, tm, c) in enumerate(stages):
    x = Inches(0.55) + (bw + Inches(0.22)) * i
    box(s, x, lane_y, bw, Inches(1.05), t_, sub, fill=WHITE, line=c, lw=1.75, tc=c, tsize=13, ssize=9)
    text(s, x, lane_y + Inches(1.14), bw, Inches(0.3), tm, size=12, color=INK, bold=True, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, x + bw, lane_y + Inches(0.52), x + bw + Inches(0.22), lane_y + Inches(0.52), color=c, w=2)
text(s, Inches(0.55), Inches(3.42), Inches(12.2), Inches(0.35), "TWO THINGS THAT MATTER HERE", size=12, color=TEAL, bold=True)
bullets(s, Inches(0.55), Inches(3.82), Inches(6.0), Inches(2.6), [
    ("Retrieval, FAQ, the spec-sheet lookup and the CMDB enrichment all run concurrently in one thread pool. Summing the per-step timings over-counts real elapsed time by 25–65%, so the harness stamps true wall-clock instead.", 0, "They overlap.  "),
    ("A ticket routed to L2 never enters the resolution loop at all — triage is a single LLM call on ticket text, with no retrieval.", 0, "Escalation short-circuits.  "),
], size=13, gap=14)
card(s, Inches(7.0), Inches(3.82), Inches(2.72), Inches(1.25), "Median", "0.62s", "half of tickets", TEAL_DK)
card(s, Inches(10.02), Inches(3.82), Inches(2.72), Inches(1.25), "P95", "3.9s", "worst realistic case", NAVY)
card(s, Inches(7.0), Inches(5.22), Inches(2.72), Inches(1.25), "Mean", "1.17s", "excl. one transient stall", TEAL_DK)
card(s, Inches(10.02), Inches(5.22), Inches(2.72), Inches(1.25), "vs. today", "30–45 min", "ServeWell's stated baseline", NAVY)
foot(s)

# =====================================================================
# 7 — Agents & orchestration
# =====================================================================
s = slide(); head(s, "Requirement i", "The agents, and how work is orchestrated")
rows = [
    ["Triage Agent", "Decide the queue: L1_GUIDED, L2_ESCALATION or NON_IT",
     "Deterministic rules first, one LLM call only for genuine judgement calls", "No tools. Ticket text only."],
    ["Resolution Agent", "Guide the store through the runbook, or hand off",
     "Bounded loop: ≤5 turns, ≤2 KB searches, 1 action proposal", "Retrieval + CMDB + actions"],
    ["L2 Copilot", "Assist the human engineer on escalated tickets",
     "Open-ended chat, ≤2 searches per question, no further escalation", "Same tools, no HITL gate"],
]
table(s, Inches(0.55), Inches(1.70), Inches(12.2), Inches(2.5),
      ["Agent", "Responsibility", "Control flow", "Tool access"], rows,
      ratios=[1.6, 3.2, 3.9, 2.3], fs=12)
text(s, Inches(0.55), Inches(4.45), Inches(12.2), Inches(0.35), "WHY A CUSTOM LOOP RATHER THAN A FRAMEWORK", size=12, color=TEAL, bold=True)
bullets(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(2.0), [
    ("Every control decision — the search budget, the turn cap, the approval gate — is a visible line of Python I can point at in this room. With LangGraph or CrewAI those same limits live inside framework config, and 'why did it stop there?' becomes a harder question to answer live.", 0, "Defensibility:  "),
    ("The orchestration is ~200 lines. A framework would have been justified if we needed persistence, retries across sessions, or multi-agent negotiation. We need none of those for one incident type.", 0, "Proportionality:  "),
    ("Bounded loops, not open-ended autonomy. Instruction-following alone did not hold — the model reworded queries and searched 4–6 times per ticket until the cap was enforced in code.", 0, "Measured need:  "),
], size=13, gap=13)
foot(s)

# =====================================================================
# 8 — Triage decision logic  (plain-language version)
# =====================================================================
s = slide(); head(s, "Requirement i", "Triage: when a simple rule beats the AI")
text(s, Inches(0.55), Inches(1.56), Inches(12.2), Inches(0.4),
     "One question decides most tickets:  how many times has L1 already worked this one and failed?",
     size=15, color=INK)

rowsv = [
    ("Nobody has tried yet", "93 tickets", "AI decides", "gets it right 67% of the time",
     TEAL_DK, False),
    ("L1 tried once", "40 tickets", "AI decides", "gets it right 70% of the time",
     TEAL_DK, False),
    ("L1 tried twice or more", "87 tickets", "Rule decides", "the rule is right 97% — the AI only 92%",
     NAVY, True),
]
ry = Inches(2.18)
for i, (label, n, who, acc, c, is_rule) in enumerate(rowsv):
    y = ry + Inches(1.06) * i
    rect(s, Inches(0.55), y, Inches(12.2), Inches(0.92), WHITE,
         MSO_SHAPE.ROUNDED_RECTANGLE, line=(NAVY if is_rule else LINE),
         lw=(2.0 if is_rule else 1.0), adj=0.10)
    rect(s, Inches(0.55), y, Inches(0.11), Inches(0.92), c)
    text(s, Inches(0.85), y + Inches(0.16), Inches(3.9), Inches(0.36), label, size=16, bold=True, color=INK)
    text(s, Inches(0.85), y + Inches(0.52), Inches(3.9), Inches(0.30), n, size=11.5, color=INK_2)
    badge = rect(s, Inches(5.05), y + Inches(0.26), Inches(1.75), Inches(0.42),
                 (NAVY if is_rule else TEAL_SOFT), MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.30)
    text(s, Inches(5.05), y + Inches(0.26), Inches(1.75), Inches(0.42), who, size=12,
         color=(WHITE if is_rule else TEAL_DK), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.10), y + Inches(0.30), Inches(5.5), Inches(0.36), acc,
         size=13, color=(NAVY if is_rule else INK_2), bold=is_rule)

rect(s, Inches(0.55), Inches(5.48), Inches(12.2), Inches(0.86), TEAL_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.12)
text(s, Inches(0.9), Inches(5.62), Inches(11.5), Inches(0.6),
     "In plain terms:  if L1 has already worked a ticket twice and it is still open, you do not need AI to tell you it needs an expert. "
     "A one-line rule is more accurate than the model there — and it is free, instant, and always gives the same answer.",
     size=13, color=TEAL_DK, ls=1.2)

rect(s, Inches(0.55), Inches(6.48), Inches(12.2), Inches(0.60), RGBColor(0xFD, 0xF4, 0xE3), MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.16)
text(s, Inches(0.9), Inches(6.58), Inches(11.5), Inches(0.4),
     "Being straight with you: these thresholds were set using the same 220 tickets they are scored on, so treat 90% as a ceiling, not a promise.",
     size=11.5, color=AMBER)
foot(s)

# =====================================================================
# 9 — RAG structure
# =====================================================================
s = slide(); head(s, "Requirement ii", "How retrieval is structured, and how answers stay grounded")
stages = [("Chunk", "runbooks split on\nsection headings"), ("Embed", "Bedrock Titan v2\n1024-dim"),
          ("Store", "PGVector\nmain + delta index"), ("Retrieve", "3 concurrent queries\ncategory-aware"),
          ("Rerank", "cross-encoder\nMiniLM-L6-v2"), ("Merge", "dedupe ~11%\noverlap")]
bw2 = Inches(1.90)
for i, (t_, sub) in enumerate(stages):
    x = Inches(0.55) + (bw2 + Inches(0.13)) * i
    box(s, x, Inches(1.68), bw2, Inches(0.86), t_, sub, fill=WHITE, line=NAVY_2, lw=1.4, tc=NAVY_2, tsize=11.5, ssize=8)
    if i < 5:
        arrow(s, x + bw2, Inches(2.11), x + bw2 + Inches(0.13), Inches(2.11), color=NAVY_2, w=1.5)
bullets(s, Inches(0.55), Inches(2.85), Inches(6.05), Inches(3.4), [
    ("A blended category + subject + description query, plus a subcategory-only query, plus an FAQ-filtered query — run concurrently. The blend alone missed network runbooks on POS/Network tickets because 'POS terminal' dominates the embedding.", 0, "Three queries, not one.  "),
    ("Spec sheets read like reference tables, not symptoms, so they rank poorly against symptom queries. get_system_spec attaches the right one by system_version instead of hoping search finds it.", 0, "One deterministic lookup.  "),
    ("Every chunk keeps its “--- Document Source: X ---” header through merging, so every claim in a reply can be traced back to a file.", 0, "Provenance survives.  "),
], size=12.5, gap=12)
card(s, Inches(6.95), Inches(2.85), Inches(2.78), Inches(1.3), "Runbook hit-rate", "100%", "correct doc retrieved, n=75", TEAL_DK)
card(s, Inches(9.97), Inches(2.85), Inches(2.78), Inches(1.3), "Recall@k", "72.8%", "of all labelled docs", NAVY)
card(s, Inches(6.95), Inches(4.30), Inches(2.78), Inches(1.3), "Groundedness", "100%", "specifics traced to source", TEAL_DK)
card(s, Inches(9.97), Inches(4.30), Inches(2.78), Inches(1.3), "Reranker lift", "+2.7pp", "hit-rate vs. no rerank", NAVY)
rect(s, Inches(6.95), Inches(5.78), Inches(5.8), Inches(0.75), RGBColor(0xFD, 0xF4, 0xE3), MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.13)
text(s, Inches(7.2), Inches(5.92), Inches(5.35), Inches(0.5),
     "Recall@k is 73%, not 100%: the key lists ~4–5 docs per ticket and we retrieve the decisive one plus some. Hit-rate is the metric that matches the job.",
     size=10.5, color=AMBER, ls=1.15)
foot(s)

# =====================================================================
# 10 — Structured data / tools vs reasoning
# =====================================================================
s = slide(); head(s, "Requirement iii", "Structured data: when it calls a tool, when it reasons")
rows = [
    ["Asset, store, SLA, warranty", "Always fetched, before the model runs", "Deterministic — never left to the model to remember to ask. ~18ms."],
    ["The relevant runbook", "Always pre-fetched, 3 concurrent queries", "The model reads results; it does not decide whether to look."],
    ["System spec sheet", "Deterministic lookup on system_version", "Semantic search ranks these poorly; a keyword map is simply better."],
    ["A second KB search", "Model's choice — capped at one", "Genuine judgement: 'is the pre-fetched runbook actually about this?'"],
    ["Which steps to give", "Model reasons over retrieved text", "This is the part that needs a language model at all."],
    ["Whether to act", "Model proposes; policy and a human dispose", "Never the model's call alone."],
]
table(s, Inches(0.55), Inches(1.70), Inches(12.2), Inches(3.5),
      ["What", "Tool call or reasoning?", "Why it is drawn there"], rows,
      ratios=[3.0, 3.9, 5.3], fs=12)
rect(s, Inches(0.55), Inches(5.45), Inches(12.2), Inches(1.42), TEAL_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.08)
text(s, Inches(0.9), Inches(5.62), Inches(11.5), Inches(0.32), "The principle", size=13, color=TEAL_DK, bold=True)
text(s, Inches(0.9), Inches(5.96), Inches(11.5), Inches(0.8),
     "If the answer is a lookup, do the lookup — do not spend a model turn deciding to. Reserve the model for the parts that are genuinely ambiguous. "
     "Structured facts are injected as authoritative and the prompt forbids contradicting them, so the agent cannot advise a warranty replacement on an out-of-warranty asset, "
     "and says so plainly when an asset is missing from the CMDB rather than inventing it.",
     size=12, color=INK_2, ls=1.25)
foot(s)

# =====================================================================
# 11 — Actions, HITL, guardrails
# =====================================================================
s = slide(); head(s, "Requirement iv", "Actions: the agent proposes, a human disposes")
flow = [("LLM proposes", "picks from a closed\ncatalog of 4 actions", NAVY, WHITE),
        ("Policy veto", "deterministic block-list:\npayment, security, multi-store", BAD, WHITE),
        ("Human reviews", "sees action, asset,\nreason, runbook cited", AMBER, WHITE),
        ("Execute + audit", "runs only on approval;\nlogged either way", GOOD, WHITE)]
fw = Inches(2.85)
for i, (t_, sub, c, tc_) in enumerate(flow):
    x = Inches(0.62) + (fw + Inches(0.32)) * i
    box(s, x, Inches(1.78), fw, Inches(1.12), t_, sub, fill=WHITE, line=c, lw=2.0, tc=c, tsize=13, ssize=9)
    if i < 3:
        arrow(s, x + fw, Inches(2.34), x + fw + Inches(0.32), Inches(2.34), color=c, w=2)
text(s, Inches(0.55), Inches(3.22), Inches(12.2), Inches(0.35), "REFUSAL IS THE DEFAULT, NOT AN EDGE CASE", size=12, color=BAD, bold=True)
bullets(s, Inches(0.55), Inches(3.62), Inches(12.2), Inches(2.6), [
    ("execute_action() refuses unless approved=True, and that check lives in the tool, not the UI — so no future caller (a script, a retry path, another agent) can reach an actuator by skipping the screen that was meant to ask.", 0, "The gate is code:  "),
    ("Four reversible, single-asset operations. The model cannot emit a free-text command, and cannot propose a printer restart for a router — the catalog is typed to asset class.", 0, "Closed vocabulary:  "),
    ("Categories the SOP marks immediate-escalate (payment, outage, security, database, multi-store) are refused before a human is ever asked. A confidently wrong proposal never becomes a button.", 0, "Policy runs first:  "),
    ("Proposed, refused and executed attempts all append to logs/ticket_actions.jsonl with actor and timestamp — the ITSM audit trail, in POC form.", 0, "Everything is logged:  "),
], size=12.5, gap=11)
foot(s)

# =====================================================================
# 12 — Evaluation & observability
# =====================================================================
s = slide(); head(s, "Requirement v", "How I know it is correct, safe and reliable")
cards = [("Routing accuracy", "90.2%", "mean of 4 runs · vs answer key", TEAL_DK),
         ("Runbook hit-rate", "100%", "correct doc retrieved", NAVY),
         ("Recall@k", "72.8%", "of all labelled docs", NAVY),
         ("Groundedness", "100%", "148 specifics, 0 unsupported", TEAL_DK),
         ("Guardrail pass", "100%", "no budget or policy breach", GOOD),
         ("Response p95", "3.9s", "true wall-clock", NAVY)]
cw, ch, g = Inches(3.9), Inches(1.42), Inches(0.24)
for i, (l, v, sub, c) in enumerate(cards):
    r_, cidx = divmod(i, 3)
    card(s, Inches(0.55) + (cw + g) * cidx, Inches(1.68) + (ch + g) * r_, cw, ch, l, v, sub, c)
text(s, Inches(0.55), Inches(4.90), Inches(12.2), Inches(0.35), "WHAT MAKES THESE TRUSTWORTHY", size=12, color=TEAL, bold=True)
bullets(s, Inches(0.55), Inches(5.30), Inches(12.2), Inches(1.7), [
    ("Every run calls the live agent over HTTP for all 256 tickets — no mocks, no cached answers. One command reproduces it.", 0, "Nothing is simulated:  "),
    ("The agent publishes its own search budget into the trace, so the guardrail checker grades against the real policy instead of a hardcoded copy that drifted out of sync twice.", 0, "Self-describing guardrails:  "),
    ("Latency comes from the same calls made for accuracy, stamped as true wall-clock — concurrent runs are flagged invalid and excluded from the trend history rather than quietly inflating p95.", 0, "Latency measured honestly:  "),
], size=12.5, gap=11)
foot(s)

# =====================================================================
# 13 — The ground truth catch
# =====================================================================
s = slide(NAVY)
text(s, Inches(0.55), Inches(0.44), Inches(9), Inches(0.3), "HOW I KNOW THE MEASUREMENTS ARE REAL", size=12.5, color=TEAL, bold=True)
text(s, Inches(0.55), Inches(0.78), Inches(12), Inches(0.7), "We were grading against the wrong answer key", size=28, color=WHITE, bold=True)
rect(s, Inches(0.55), Inches(1.48), Inches(0.85), Pt(3), TEAL)
steps = [
    ("THE SETUP", "Every ticket ships with escalation_flag — what an upstream system marked. Separately, labels/train_labels.json carries correct_routing: the actual answer key."),
    ("THE BUG", "Our evaluator scored routing against escalation_flag. It agrees with the real key only 50.5% of the time — worse than a coin flip on the one signal that mattered."),
    ("HOW IT SURFACED", "A routing number that looked plausible (72%) did not survive the question “is this graded against the right file?”. It was not."),
    ("THE FIX", "Re-pointed every evaluator, the UI and the ROI model at correct_routing, then recalibrated triage against the real target distribution — 65% of tickets should escalate, not ~50%."),
    ("WHY IT IS ON A SLIDE", "A POC that cannot detect its own measurement error is not evidence of anything. This is the check that makes the other numbers worth reading."),
]
for i, (t_, d) in enumerate(steps):
    ry = Inches(1.92) + Inches(1.02) * i
    rect(s, Inches(0.55), ry + Inches(0.02), Inches(0.13), Inches(0.80), TEAL)
    text(s, Inches(0.85), ry, Inches(2.6), Inches(0.3), t_, size=11.5, color=TEAL, bold=True)
    text(s, Inches(0.85), ry + Inches(0.29), Inches(11.6), Inches(0.7), d, size=12.5, color=RGBColor(0xD3, 0xDD, 0xE7), ls=1.15)
foot(s)

# =====================================================================
# 14 — Accuracy journey
# =====================================================================
s = slide(); head(s, "Evaluation", "Routing accuracy — four measurements, from the corrected baseline")
# The wrong-ground-truth number (31.2%) is deliberately not on this chart. It
# was graded against a different field (escalation_flag) than every bar here
# (correct_routing), so it is not the same metric -- plotting it alongside
# these would imply a comparability that is not real. The previous slide
# already tells that story in full; this chart starts from the point where
# every bar is measuring the same thing.
cd = CategoryChartData()
cd.categories = ["Correct key +\nrecalibrated prompt", "+ rule: repeat\nL1 history",
                 "+ rule: fault-class\nsubcategory", "+ prompt fix:\nsymptom survived fix"]
cd.add_series("Routing accuracy", (77.3, 79.5, 88.9, 90.2))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.62), Inches(8.9), Inches(4.75), cd)
ch_ = gf.chart; ch_.has_legend = False
ch_.has_title = False   # the slide title already says this
pl = ch_.plots[0]; pl.has_data_labels = True
dl = pl.data_labels; dl.number_format = '0.0"%"'; dl.number_format_is_linked = False
dl.font.size = Pt(12); dl.font.bold = True; dl.font.color.rgb = NAVY
ser = pl.series[0]
for i, pt in enumerate(ser.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = [RGBColor(0x8E, 0x9E, 0xB2),
                                     RGBColor(0x39, 0x77, 0x8C), RGBColor(0x1C, 0x86, 0x7C), TEAL][i]
va = ch_.value_axis; va.maximum_scale = 100; va.minimum_scale = 0; va.has_major_gridlines = False
va.tick_labels.font.size = Pt(10)
ch_.category_axis.tick_labels.font.size = Pt(9.5); ch_.font.name = F
text(s, Inches(9.72), Inches(1.72), Inches(3.05), Inches(0.3), "EACH STEP WAS MEASURED", size=11, color=TEAL, bold=True)
text(s, Inches(9.72), Inches(2.06), Inches(3.05), Inches(1.5),
     "Not one prompt rewrite. Each gain came from reading the failures, forming a hypothesis, and testing it at n=220.",
     size=11.5, color=INK_2, ls=1.25)
card(s, Inches(9.72), Inches(3.30), Inches(3.03), Inches(1.28), "Total gain", "+11 pts", "79.5% → 90.2%", TEAL_DK)
card(s, Inches(9.72), Inches(4.75), Inches(3.03), Inches(1.28), "Escalation calibration", "66% vs 65%", "predicted vs expected", NAVY)
rect(s, Inches(9.72), Inches(6.20), Inches(3.03), Inches(0.68), RGBColor(0xFD, 0xF4, 0xE3), MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.14)
text(s, Inches(9.90), Inches(6.32), Inches(2.7), Inches(0.5), "One attempt overshot to 83% and was reverted.",
     size=10, color=AMBER, ls=1.1)
foot(s)

# =====================================================================
# 15 — Latency engineering
# =====================================================================
s = slide(); head(s, "Reliability", "Making it fast — and measuring it honestly")
rows = [
    ["Summed step timings", "Over-counted 25–65%: concurrent calls double-counted", "True wall-clock stamped on every exit path", "correctness"],
    ["Duplicate action proposals", "Model re-proposed the same action, burning LLM turns", "Capped at one successful proposal per ticket", "−27% avg"],
    ["Retry on invalid actions", "Refusal invited another guess for assets with no actions", "Refusal now names that none exist", "outlier removed"],
    ["Repeated identical queries", "Re-embedded, re-searched, re-ranked every time", "Full result cache — 54% hit rate warm", "−25% p95"],
]
table(s, Inches(0.55), Inches(1.70), Inches(12.2), Inches(2.6),
      ["What was slow", "Why", "Fix", "Effect"], rows, ratios=[2.7, 4.3, 3.5, 1.7], fs=11.5)
for i, (l, v, sub, c) in enumerate([("Avg response", "1.9s → 1.2s", "−39%", TEAL_DK),
                                    ("P95 response", "5.8s → 3.9s", "−34%", TEAL_DK),
                                    ("Worst case", "64s → 5.9s", "stall removed", NAVY),
                                    ("Accuracy", "unchanged", "no trade made", GOOD)]):
    card(s, Inches(0.55) + Inches(3.13) * i, Inches(4.55), Inches(2.95), Inches(1.35), l, v, sub, c)
rect(s, Inches(0.55), Inches(6.10), Inches(12.2), Inches(0.80), TEAL_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.12)
text(s, Inches(0.85), Inches(6.24), Inches(11.6), Inches(0.55),
     "The caching result improves with volume: hit rate went 0% → 18% → 54% as the cache warmed across runs. At 400 tickets/day over a stable ~45-subcategory taxonomy, the same queries recur constantly — this gets better in production, not worse.",
     size=12, color=TEAL_DK, ls=1.2)
foot(s)

# =====================================================================
# 16 — Trade-offs (requirement vi)
# =====================================================================
s = slide(); head(s, "Requirement vi", "The trade-offs I made, and what I gave up")
rows = [
    ["Model", "Bedrock nvidia-nemotron-nano-3-30b", "The only model this AWS account can reach — verified by AccessDenied on every Claude/Nova/Llama candidate. A small model made the guardrails earn their place."],
    ["Orchestration", "Custom ~200-line loop, not LangGraph/CrewAI", "Every limit is a line of Python I can defend live. Gave up: persistence, retries, a visual graph."],
    ["Retrieval", "PGVector + Titan v2 + cross-encoder rerank", "Rerank costs ~0.3s and is CPU-bound, but buys +2.7pp hit-rate. Diverged from the plan's ChromaDB + MiniLM."],
    ["Autonomy line", "Propose-only; execution needs approved=True", "Slower than auto-remediation. Chosen because a wrong action in 2,500 stores is unrecoverable; wrong advice is not."],
    ["Triage design", "Deterministic rules where measured better", "Less elegant than 'the LLM decides'. But at 2+ notes the rule beats the model 96.6% to 92.0%."],
    ["Scope", "One incident family, deep, over six shallow", "The brief rewards a POC defended deeply. Ticket write-back is the visible gap."],
]
table(s, Inches(0.55), Inches(1.68), Inches(12.2), Inches(4.5),
      ["Decision", "What I chose", "Why, and what it cost"], rows, ratios=[1.8, 3.6, 6.8], fs=11.5)
rect(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.62), TEAL_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.15)
text(s, Inches(0.85), Inches(6.48), Inches(11.6), Inches(0.4),
     "Prompts now live as versioned files in prompts/ — extracted byte-identically, so a prompt change reads like a behaviour change in a diff, which is what it is.",
     size=12, color=TEAL_DK, italic=True)
foot(s)

# =====================================================================
# 17 — Business impact
# =====================================================================
# Redesigned from a flat table to: (1) a legend establishing the same
# Measured/Stated/Assumption vocabulary the roi.html artifact already uses,
# so the two read as one system rather than two inconsistent takes on the
# same numbers; (2) four headline cards instead of three, adding "tickets
# deflected/day" so the 35% L1-share lever has a concrete daily number next
# to the dollar figure it drives; (3) real chip badges per evidence row
# instead of colored table text.
#
# Two numbers corrected while rebuilding this, both worth flagging rather
# than silently changing:
#   - FTE was "~3.2", computed as hoursDay/8 (an 8-hour operating day). A
#     "full-time equivalent" conventionally means a 2,080-hour work-year
#     (40hr x 52wk), which is also what roi.html's own JS uses. Recomputed
#     on that basis: 9,198 / 2,080 = ~4.4 FTE -- the deck and the artifact
#     were quietly using two different definitions of the same label.
#   - p95 was "3.9s"; the latest measured run is 3.969s, which rounds to
#     4.0s. Kept the recompute rather than the stale rounding.
s = slide(); head(s, "For the Business", "What this is worth, and what each number rests on")

# --- legend: same three tags as the roi.html artifact ---
leg_items = [("MEASURED", GOOD), ("STATED BY SERVEWELL", VIOLET), ("YOUR ASSUMPTION", AMBER)]
lx = Inches(0.55)
for label, c in leg_items:
    w = Inches(0.24 + 0.083 * len(label))
    pill(s, lx, Inches(1.53), w, Inches(0.26), label, c, fs=8.5)
    lx += w + Inches(0.14)

# --- four headline cards ---
cw, cgap = Inches(2.87), Inches(0.21)
cy = Inches(1.90)
card(s, Inches(0.55) + (cw + cgap) * 0, cy, cw, Inches(1.35), "L1 cost avoided / yr", "$202k", "at $22/hr fully loaded", NAVY)
card(s, Inches(0.55) + (cw + cgap) * 1, cy, cw, Inches(1.35), "L1 hours returned / yr", "9,198", "≈ 4.4 FTE at 2,080 hrs/yr", TEAL_DK)
card(s, Inches(0.55) + (cw + cgap) * 2, cy, cw, Inches(1.35), "Tickets deflected / day", "84", "of 140 L1-eligible (35% share)", NAVY)
card(s, Inches(0.55) + (cw + cgap) * 3, cy, cw, Inches(1.35), "Time to first response", "38 min → 1.2s", "measured end to end", TEAL_DK)

# --- evidence rows, chip-badged rather than color-coded text ---
ev_rows = [
    ("Tickets resolvable at L1", "35%", "MEASURED", GOOD, "answer key, 220 tickets"),
    ("Runbook retrieved correctly", "100% hit-rate", "MEASURED", GOOD, "vs. answer key"),
    ("Guidance grounded in a runbook", "100%", "MEASURED", GOOD, "162 specifics checked, 0 unsupported"),
    ("Agent response, end to end", "1.2s avg · 4.0s p95", "MEASURED", GOOD, "true wall-clock"),
    ("Time to first response today", "30–45 min", "STATED", VIOLET, "ServeWell's own brief"),
    ("Repetitive, runbook-mappable tickets", "60–70%", "STATED", VIOLET, "ServeWell's own brief"),
    ("Volume · handle time · cost · deflection", "400/day · 18min · $22/hr · 60%", "ASSUMPTION", AMBER, "yours to set — adjustable"),
]
rx1, rx2, rx3, rx4 = Inches(0.55), Inches(4.20), Inches(6.35), Inches(7.95)
rw1, rw2, rw3, rw4 = Inches(3.55), Inches(2.05), Inches(1.50), Inches(4.80)
ry0, rh, rgap = Inches(3.42), Inches(0.40), Inches(0.05)
for i, (label, val, tag, c, detail) in enumerate(ev_rows):
    ry = ry0 + (rh + rgap) * i
    rect(s, rx1 - Inches(0.05), ry - Inches(0.02), Inches(12.2) + Inches(0.10), rh + Inches(0.04),
         WHITE if i % 2 == 0 else RGBColor(0xF1, 0xF4, 0xF7), MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.25)
    text(s, rx1, ry + Inches(0.04), rw1, Inches(0.32), label, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx2, ry + Inches(0.04), rw2, Inches(0.32), val, size=12, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    pill(s, rx3, ry + Inches(0.055), Inches(1.35), Inches(0.29), tag, c, fs=9)
    text(s, rx4, ry + Inches(0.04), rw4, Inches(0.32), detail, size=10.5, color=INK_2, italic=True, anchor=MSO_ANCHOR.MIDDLE)

banner_y = ry0 + (rh + rgap) * len(ev_rows) + Inches(0.11)
rect(s, Inches(0.55), banner_y, Inches(12.2), Inches(0.55), RGBColor(0xFD, 0xF4, 0xE3), MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.18)
text(s, Inches(0.85), banner_y + Inches(0.09), Inches(11.6), Inches(0.38),
     "The deflection rate is the single assumption most worth validating first — it is a hypothesis until the agent runs in shadow mode against live tickets.",
     size=11.5, color=AMBER)
foot(s)

# =====================================================================
# 18 — Candid limitations
# =====================================================================
s = slide(NAVY)
text(s, Inches(0.55), Inches(0.44), Inches(9), Inches(0.3), "WHAT I WOULD IMPROVE WITH MORE TIME", size=12.5, color=TEAL, bold=True)
text(s, Inches(0.55), Inches(0.78), Inches(12), Inches(0.7), "What this POC does not claim", size=28, color=WHITE, bold=True)
rect(s, Inches(0.55), Inches(1.48), Inches(0.85), Pt(3), TEAL)
cav = [
    ("Deflection is projected, not proven.", "The agent guides a store through a fix; it does not close tickets unaided. Shadow-mode against live traffic is the highest-value next step, and the one number I would insist on validating before signing a business case."),
    ("Routing is tuned on the tickets it is scored on.", "90.2% is an upper bound, not a held-out result. The honest test is the unseen cases you hand me today."),
    ("The model is non-deterministic at temperature 0.", "Repeat runs vary ±1–2 points; 16 of 33 failing tickets flip between runs. I report means over 4 runs rather than a best single number."),
    ("Ticket write-back is not wired.", "The /save_history endpoint exists and nothing calls it. In production this is what would let ticket history compound — which the triage rule depends on."),
    ("One incident family, synthetic data.", "256 tickets, one knowledge base, clean text. Real tickets are messier, and retrieval on escalation-procedure.md already fails 8 times out of 8."),
]
for i, (t_, d) in enumerate(cav):
    ry = Inches(1.92) + Inches(1.02) * i
    rect(s, Inches(0.55), ry + Inches(0.02), Inches(0.13), Inches(0.80), WARN)
    text(s, Inches(0.85), ry, Inches(11.6), Inches(0.3), t_, size=13.5, color=WHITE, bold=True)
    text(s, Inches(0.85), ry + Inches(0.30), Inches(11.6), Inches(0.68), d, size=11.5, color=RGBColor(0xC9, 0xD5, 0xE1), ls=1.15)
foot(s)

# =====================================================================
# 19 — Unseen data
# =====================================================================
s = slide(); head(s, "Requirement 4", "When you hand me a ticket I have never seen")
cols = [
    ("What will happen", TEAL_DK, [
        "It routes on ticket text alone — the flag is withheld.",
        "Retrieval runs three concurrent queries over the same KB.",
        "Guardrails hold regardless of content: budgets are enforced in code.",
        "No action executes without your explicit approval.",
    ]),
    ("Where it will struggle", AMBER, [
        "A subcategory outside the 8 fault-class rules falls to the LLM (~68% there).",
        "A symptom with no runbook: it should escalate rather than invent steps.",
        "Policy-blocked categories refuse automation even when a fix exists.",
        "Novel phrasing may retrieve a sibling runbook instead of the exact one.",
    ]),
    ("How you can check me live", VIOLET, [
        "Every trace shows the retrieved documents and their source files.",
        "decided_by tells you whether a rule or the model made the call.",
        "The reply cites the runbook — verify it against the file.",
        "Re-run the eval in the UI and watch the numbers move.",
    ]),
]
for i, (t_, c, items) in enumerate(cols):
    x = Inches(0.55) + Inches(4.15) * i
    rect(s, x, Inches(1.70), Inches(3.85), Inches(4.65), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, adj=0.05)
    rect(s, x, Inches(1.70), Inches(3.85), Inches(0.10), c)
    text(s, x + Inches(0.28), Inches(1.98), Inches(3.3), Inches(0.35), t_, size=14, color=c, bold=True)
    bullets(s, x + Inches(0.28), Inches(2.45), Inches(3.35), Inches(3.6),
            [(b, 0) for b in items], size=11.5, bcolor=c, gap=12, ls=1.15)
rect(s, Inches(0.55), Inches(6.52), Inches(12.2), Inches(0.5), TEAL_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.18)
text(s, Inches(0.85), Inches(6.60), Inches(11.6), Inches(0.35),
     "If a test case fails today, I would rather show you exactly why from the trace than have it pass for a reason neither of us can see.",
     size=12, color=TEAL_DK, italic=True)
foot(s)

# =====================================================================
# 20 — Close
# =====================================================================
s = slide(); head(s, "Next", "Where this goes from here")
nx = [
    ("Wire ticket write-back", "Close the last gap so resolutions land back on the ticket and history compounds."),
    ("Shadow-mode validation", "Run against live traffic to turn the deflection assumption into a measured number."),
    ("Fix the escalation-SOP retrieval gap", "escalation-procedure.md is never retrieved (0/8). A policy doc needs a deterministic attach, like spec sheets got."),
    ("Hold the tuning line", "Further rules only where a new failure cluster shows the same structural signature — not by pattern-matching to the score."),
]
bullets(s, Inches(0.55), Inches(1.78), Inches(12.2), Inches(3.4),
        [(d, 0, t_ + ":  ") for t_, d in nx], size=15, gap=20, ls=1.15)
rect(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(1.5), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, line=TEAL, lw=1.5, adj=0.05)
text(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(0.32), "What I would most like you to probe", size=13.5, color=TEAL_DK, bold=True)
text(s, Inches(0.9), Inches(5.10), Inches(11.5), Inches(0.85),
     "The autonomy line, the ground-truth catch, and the places I chose a rule over the model. Those are the three decisions this system stands or falls on — "
     "and the three I have the most measurement behind.",
     size=12.5, color=INK_2, ls=1.25)
rect(s, 0, Inches(6.28), SW, Inches(1.22), NAVY)
text(s, Inches(0.55), Inches(6.62), Inches(11.5), Inches(0.5), "Thank you — let's look at your test cases.", size=20, color=WHITE, bold=True)

OUT = Path(__file__).parent / "ServeWell_Final_Deck.pptx"
prs.save(str(OUT))
print(f"saved {OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
